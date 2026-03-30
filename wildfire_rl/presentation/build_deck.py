"""
build_deck.py
Generate a Gartner-style presentation deck from experiment results.

Usage:
    python presentation/build_deck.py

Output: presentation/Wildfire_RL_Presentation.pptx
"""

import sys
import pickle
from pathlib import Path
from io import BytesIO

import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, str(Path(__file__).parent.parent))

RESULTS_DIR = Path(__file__).parent.parent / "results"
OUTPUT_DIR = Path(__file__).parent
OUTPUT_PATH = OUTPUT_DIR / "Wildfire_RL_Presentation.pptx"

# =====================================================================
# Gartner-Style Color Palette
# =====================================================================
NAVY       = RGBColor(0x0B, 0x0B, 0x2B)   # deep navy background
DARK_BLUE  = RGBColor(0x12, 0x1A, 0x3E)   # slide body bg
ACCENT     = RGBColor(0x00, 0x96, 0xD6)   # bright blue accent
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)   # orange for PPO / highlights
GREEN      = RGBColor(0x27, 0xAE, 0x60)   # green for greedy
GRAY       = RGBColor(0x95, 0xA5, 0xA6)   # gray for random
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DIM_WHITE  = RGBColor(0xAA, 0xAA, 0xAA)
RED        = RGBColor(0xE7, 0x4C, 0x3C)


# =====================================================================
# Helper Functions
# =====================================================================

def set_slide_bg(slide, color=NAVY):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add Gartner-style title bar at top of slide."""
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Title
    add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
                title_text, font_size=28, color=WHITE, bold=True)

    # Subtitle
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.4),
                    subtitle_text, font_size=14, color=DIM_WHITE)


def add_image(slide, image_path, left, top, width=None, height=None):
    """Add an image to a slide."""
    kwargs = {"image_file": str(image_path), "left": left, "top": top}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    slide.shapes.add_picture(**kwargs)


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, spacing=Pt(6)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0


def add_stat_box(slide, left, top, value, label, value_color=ACCENT):
    """Add a large stat callout box."""
    # Value
    add_textbox(slide, left, top, Inches(2.5), Inches(0.7),
                str(value), font_size=36, color=value_color, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, left, top + Inches(0.6), Inches(2.5), Inches(0.4),
                label, font_size=12, color=DIM_WHITE,
                alignment=PP_ALIGN.CENTER)


def load_results():
    path = RESULTS_DIR / "experiment_results.pkl"
    with open(path, "rb") as f:
        return pickle.load(f)


# =====================================================================
# Slide Builders
# =====================================================================

def build_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Large title
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
                "Wildfire Suppression via\nReinforcement Learning",
                font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    # Accent line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(3.6), Inches(3), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Subtitle
    add_textbox(slide, Inches(0.8), Inches(3.9), Inches(10), Inches(0.8),
                "Comparing RL Algorithms on a Forest Fire Cellular Automaton Environment",
                font_size=18, color=DIM_WHITE)

    # Course info
    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(10), Inches(0.5),
                "MMAI-845  |  Team Union", font_size=14, color=GRAY)


def build_problem_slide(prs):
    """Slide 2: Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Why Wildfire Suppression?", "The problem and why RL is the right approach")

    items = [
        "Wildfire containment is a time-critical, spatially constrained resource routing problem",
        "Fire spreads stochastically — every timestep of inaction expands the fire front",
        "Incident commanders must make real-time routing decisions under extreme cognitive load",
        "RL agents can learn adaptive suppression policies through thousands of simulated episodes",
        "This project compares 4 agent strategies on a standardized forest fire environment",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(7), Inches(4.5),
                    items, font_size=18, color=LIGHT_GRAY)

    # Key stat callout
    add_stat_box(slide, Inches(9), Inches(1.8), "2,200+", "Wildfires in BC (2023)")
    add_stat_box(slide, Inches(9), Inches(3.2), "$720M", "Suppression Costs")
    add_stat_box(slide, Inches(9), Inches(4.6), "2.84M ha", "Area Burned")


def build_env_slide(prs):
    """Slide 3: Environment Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Environment: ForestFireHelicopter5x5-v1",
                  "Pre-defined Gymnasium environment from gym-cellular-automata")

    # Left: description
    items = [
        "5x5 grid  |  Drossel-Schwabl cellular automaton",
        "Cell states: Empty (0), Tree (1), Fire (2)",
        "Fire spreads to adjacent trees, burns out after 1 step",
        "Random lightning strikes (p=0.033) start new fires",
        "Trees regrow randomly (p=0.333)",
        "Helicopter extinguishes fire by flying over burning cells",
        "9 actions: 8 directions + stay",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(4),
                    items, font_size=16)

    # Right: grid snapshot image
    img_path = RESULTS_DIR / "08_grid_snapshots.png"
    if img_path.exists():
        add_image(slide, img_path, Inches(6.5), Inches(1.3), width=Inches(6.3))

    # Reward formula
    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(0.5),
                "Reward = (trees \u2212 fires) / 25  +  extinguish bonus  +  proximity bonus",
                font_size=14, color=ACCENT, bold=True)


def build_algo_overview(prs):
    """Slide 4: Algorithm Overview — 2x2 matrix"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Four Agents, Four Strategies",
                  "From trivial baseline to deep reinforcement learning")

    # 2x2 grid of agent descriptions
    agents = [
        ("Random", GRAY, "Trivial Baseline",
         "Uniformly random actions.\nEstablishes the performance floor.\nNo intelligence whatsoever."),
        ("Greedy (BFS)", GREEN, "Heuristic Baseline",
         "BFS to nearest fire, move toward it.\nPerfect grid information access.\nSmart but not learning."),
        ("DQN", ACCENT, "Deep RL \u2014 Value-Based",
         "Learns Q-values via replay buffer.\nOff-policy: reuses past experience.\nNeural net predicts action values."),
        ("PPO", ORANGE, "Deep RL \u2014 Policy Gradient",
         "Learns action probabilities directly.\nOn-policy with clipped updates.\nEntropy bonus for exploration."),
    ]

    positions = [
        (Inches(0.6), Inches(1.5)),
        (Inches(6.6), Inches(1.5)),
        (Inches(0.6), Inches(4.2)),
        (Inches(6.6), Inches(4.2)),
    ]

    for (name, color, subtitle, desc), (left, top) in zip(agents, positions):
        # Box background
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(5.8), Inches(2.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BLUE
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Agent name
        add_textbox(slide, left + Inches(0.3), top + Inches(0.15),
                    Inches(5), Inches(0.4),
                    name, font_size=20, color=color, bold=True)

        # Subtitle
        add_textbox(slide, left + Inches(0.3), top + Inches(0.55),
                    Inches(5), Inches(0.3),
                    subtitle, font_size=12, color=DIM_WHITE)

        # Description
        add_textbox(slide, left + Inches(0.3), top + Inches(0.9),
                    Inches(5.2), Inches(1.2),
                    desc, font_size=13, color=LIGHT_GRAY)


def build_dqn_deep_dive(prs):
    """Slide 5: DQN Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "DQN — Deep Q-Network", "Off-policy, value-based deep reinforcement learning")

    items = [
        "Neural network approximates Q(s, a) — the expected reward of taking action a in state s",
        "Replay buffer (100K experiences) enables learning from past transitions",
        "Epsilon-greedy exploration: starts random, gradually becomes greedy",
        "Target network updated every 250 steps for training stability",
        "Network: 30 inputs \u2192 128 \u2192 128 \u2192 9 outputs (one Q-value per action)",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(3.5),
                    items, font_size=16)

    # Hyperparameter table
    params = [
        ("Learning Rate", "5e-4"),
        ("Buffer Size", "100,000"),
        ("Batch Size", "128"),
        ("Gamma (\u03b3)", "0.995"),
        ("Exploration", "50% \u2192 5%"),
        ("Training Steps", "500,000"),
    ]

    add_textbox(slide, Inches(8), Inches(1.4), Inches(4), Inches(0.4),
                "Key Hyperparameters", font_size=16, color=ACCENT, bold=True)

    for i, (param, value) in enumerate(params):
        y = Inches(1.9) + Inches(0.4) * i
        add_textbox(slide, Inches(8), y, Inches(2.5), Inches(0.35),
                    param, font_size=13, color=DIM_WHITE)
        add_textbox(slide, Inches(10.5), y, Inches(2), Inches(0.35),
                    value, font_size=13, color=WHITE, bold=True)


def build_ppo_deep_dive(prs):
    """Slide 6: PPO Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "PPO — Proximal Policy Optimization", "On-policy, policy gradient deep reinforcement learning")

    items = [
        "Learns a policy \u03c0(a|s) — the probability distribution over actions given state",
        "On-policy: collects fresh trajectories (1,024 steps), then updates",
        "Clipped surrogate objective prevents destructive large policy updates",
        "Entropy bonus (0.02) encourages continued exploration",
        "Network: 30 inputs \u2192 128 \u2192 128 \u2192 9 action probabilities + value estimate",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(3.5),
                    items, font_size=16)

    params = [
        ("Learning Rate", "3e-4"),
        ("Rollout Steps", "1,024"),
        ("Update Epochs", "15"),
        ("Gamma (\u03b3)", "0.995"),
        ("Clip Range", "0.2"),
        ("Entropy Coeff", "0.02"),
    ]

    add_textbox(slide, Inches(8), Inches(1.4), Inches(4), Inches(0.4),
                "Key Hyperparameters", font_size=16, color=ORANGE, bold=True)

    for i, (param, value) in enumerate(params):
        y = Inches(1.9) + Inches(0.4) * i
        add_textbox(slide, Inches(8), y, Inches(2.5), Inches(0.35),
                    param, font_size=13, color=DIM_WHITE)
        add_textbox(slide, Inches(10.5), y, Inches(2), Inches(0.35),
                    value, font_size=13, color=WHITE, bold=True)


def build_training_curves_slide(prs):
    """Slide 7: Training Curves"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Training Curves — Learning Progression",
                  "Both agents improve significantly over 500K timesteps (~2,500 episodes)")

    img_path = RESULTS_DIR / "01_training_curves.png"
    if img_path.exists():
        add_image(slide, img_path, Inches(0.5), Inches(1.3), width=Inches(8.5))

    # Key insight callout
    add_textbox(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(3),
                "Key Insight\n\n"
                "Both DQN and PPO show clear learning progression. "
                "DQN improves steadily via replay buffer. "
                "PPO shows more variance due to on-policy rollouts but converges to similar performance.",
                font_size=13, color=LIGHT_GRAY)


def build_results_success_slide(prs, results):
    """Slide 8: Success Rate Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Results — Fire Suppression Success Rate",
                  "Percentage of episodes where all fire was extinguished")

    img_path = RESULTS_DIR / "02_success_rate.png"
    if img_path.exists():
        add_image(slide, img_path, Inches(0.3), Inches(1.3), width=Inches(7.5))

    # Stat callouts
    if "Greedy" in results:
        add_stat_box(slide, Inches(8.5), Inches(1.8),
                     f"{results['Greedy']['success_rate']:.0f}%", "Greedy (Best)", GREEN)
    if "PPO" in results:
        add_stat_box(slide, Inches(8.5), Inches(3.2),
                     f"{results['PPO']['success_rate']:.0f}%", "PPO", ORANGE)
    if "DQN" in results:
        add_stat_box(slide, Inches(8.5), Inches(4.6),
                     f"{results['DQN']['success_rate']:.0f}%", "DQN", ACCENT)
    if "Random" in results:
        add_stat_box(slide, Inches(11), Inches(1.8),
                     f"{results['Random']['success_rate']:.0f}%", "Random (Floor)", GRAY)


def build_results_reward_slide(prs, results):
    """Slide 9: Mean Reward Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Results — Mean Episode Reward",
                  "100-episode evaluation with error bars (\u00b11 std)")

    img_path = RESULTS_DIR / "03_mean_reward.png"
    if img_path.exists():
        add_image(slide, img_path, Inches(0.3), Inches(1.3), width=Inches(7.5))

    items = [
        f"Greedy:  {results.get('Greedy', {}).get('mean_reward', 0):.1f} (benchmark)",
        f"DQN:     {results.get('DQN', {}).get('mean_reward', 0):.1f} (+55% vs Random)",
        f"PPO:     {results.get('PPO', {}).get('mean_reward', 0):.1f} (+52% vs Random)",
        f"Random:  {results.get('Random', {}).get('mean_reward', 0):.1f} (floor)",
    ]
    add_bullet_list(slide, Inches(8.5), Inches(1.8), Inches(4), Inches(3),
                    items, font_size=14)


def build_behavioral_slide(prs):
    """Slide 10: Position Heatmaps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Behavioral Analysis — Where Does Each Agent Go?",
                  "Position frequency heatmaps across 100 evaluation episodes")

    img_path = RESULTS_DIR / "04_position_heatmaps.png"
    if img_path.exists():
        add_image(slide, img_path, Inches(0.3), Inches(1.3), width=Inches(12.5))

    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(12), Inches(0.5),
                "Random spreads uniformly  |  Greedy concentrates on fire zones  |  "
                "RL agents develop spatial preferences through learning",
                font_size=13, color=DIM_WHITE, alignment=PP_ALIGN.CENTER)


def build_fire_speed_slide(prs):
    """Slide 11: Fire Suppression Speed"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Fire Suppression Speed",
                  "Average fire cells remaining over the episode timeline")

    img_path = RESULTS_DIR / "05_fire_over_time.png"
    if img_path.exists():
        add_image(slide, img_path, Inches(0.5), Inches(1.3), width=Inches(8.5))

    add_textbox(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(3),
                "Key Insight\n\n"
                "Greedy suppresses fire fastest due to direct BFS targeting. "
                "RL agents show intermediate performance — they learn to move toward fire "
                "but without perfect grid access, they are less efficient.",
                font_size=13, color=LIGHT_GRAY)


def build_hyperparam_slide(prs):
    """Slide 12: Hyperparameter Sensitivity"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Hyperparameter Sensitivity — Learning Rate",
                  "How learning rate affects final performance (100K timesteps, last 50 episodes)")

    img_path = RESULTS_DIR / "07_hyperparam_sweep.png"
    if img_path.exists():
        add_image(slide, img_path, Inches(0.3), Inches(1.3), width=Inches(12.5))

    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(12), Inches(0.5),
                "Both algorithms are sensitive to learning rate  |  "
                "Too low = slow convergence  |  Too high = instability",
                font_size=13, color=DIM_WHITE, alignment=PP_ALIGN.CENTER)


def build_reward_shaping_slide(prs):
    """Slide 13: Reward Shaping — Critical Design Decision"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Reward Shaping — A Critical Design Decision",
                  "Why the base reward wasn't enough, and how we fixed it")

    items_left = [
        "The Problem:",
        "Base reward = (trees \u2212 fires) / 25 per step",
        "Dominated by stochastic fire spawning (random lightning)",
        "Agent's actions had tiny marginal effect on reward",
        "Result: DQN and PPO failed to learn (= Random performance)",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(3),
                    items_left, font_size=15)

    items_right = [
        "The Fix:",
        "+2.0 bonus for extinguishing a fire cell (direct action reward)",
        "+0.5 proximity bonus for being close to fire (guidance signal)",
        "Result: 2x reward gap between Greedy and Random",
        "DQN and PPO now learn clearly (+55% above Random)",
    ]
    add_bullet_list(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3),
                    items_right, font_size=15)

    add_textbox(slide, Inches(0.8), Inches(5.2), Inches(12), Inches(1),
                "Lesson: In highly stochastic environments, the agent needs a reward signal "
                "that directly reflects its actions, not just the global system state. "
                "Reward engineering is as important as algorithm selection.",
                font_size=15, color=ACCENT, bold=True)


def build_training_time_slide(prs):
    """Slide 14: Training Efficiency"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Training Efficiency",
                  "Wall-clock time for 500K timesteps on a single machine")

    img_path = RESULTS_DIR / "06_training_time.png"
    if img_path.exists():
        add_image(slide, img_path, Inches(1), Inches(1.3), width=Inches(5.5))

    items = [
        "DQN: Off-policy replay is computationally cheaper per step",
        "PPO: On-policy rollouts + multiple gradient epochs add overhead",
        "Both trained on a single CPU — no GPU required for this environment",
        "The 5x5 grid keeps state space small enough for efficient training",
    ]
    add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(3),
                    items, font_size=15)


def build_takeaways_slide(prs, results):
    """Slide 15: Key Takeaways & Future Work"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Key Takeaways")

    takeaways = [
        "1.  RL agents can learn fire suppression strategies from scratch — both DQN and PPO "
        "significantly outperform random baselines",
        "2.  Domain heuristics (Greedy BFS) outperform general-purpose RL on small, fully-observable "
        "problems — but RL scales to scenarios where heuristics break down",
        "3.  Reward shaping is critical in stochastic environments — the agent needs a signal "
        "that reflects its actions, not just the system state",
        "4.  PPO shows better success rate than DQN (19% vs 13%) despite similar mean reward — "
        "on-policy learning may produce more consistent suppression behavior",
        "5.  Hyperparameter sensitivity is real — both algorithms are sensitive to learning rate, "
        "reinforcing the need for systematic tuning",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(3.5),
                    takeaways, font_size=16, spacing=Pt(10))

    # Future work
    add_textbox(slide, Inches(0.6), Inches(5.2), Inches(4), Inches(0.4),
                "Future Work", font_size=20, color=ACCENT, bold=True)

    future = [
        "Scale to larger grids (16x16, 32x32) where heuristics fail",
        "Add wind dynamics and terrain features",
        "Multi-agent coordination (multiple helicopters)",
        "Transfer learning across environment configurations",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(2),
                    future, font_size=14, color=DIM_WHITE)


def build_thank_you_slide(prs):
    """Slide 16: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0), Inches(2.5), Inches(13.33), Inches(1),
                "Thank You", font_size=44, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(3.6), Inches(2.33), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_textbox(slide, Inches(0), Inches(4.0), Inches(13.33), Inches(0.5),
                "Questions?", font_size=22, color=DIM_WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0), Inches(5.5), Inches(13.33), Inches(0.5),
                "Environment: gym-cellular-automata  |  Framework: Stable-Baselines3  |  "
                "Code: github.com/team-union/wildfire-rl",
                font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# Main Builder
# =====================================================================

def build_presentation():
    print("Building presentation deck...\n")

    results = load_results()

    prs = Presentation()
    prs.slide_width = Inches(13.33)   # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    build_title_slide(prs)              # 1
    build_problem_slide(prs)            # 2
    build_env_slide(prs)                # 3
    build_algo_overview(prs)            # 4
    build_dqn_deep_dive(prs)            # 5
    build_ppo_deep_dive(prs)            # 6
    build_training_curves_slide(prs)    # 7
    build_results_success_slide(prs, results)   # 8
    build_results_reward_slide(prs, results)    # 9
    build_behavioral_slide(prs)         # 10
    build_fire_speed_slide(prs)         # 11
    build_hyperparam_slide(prs)         # 12
    build_reward_shaping_slide(prs)     # 13
    build_training_time_slide(prs)      # 14
    build_takeaways_slide(prs, results) # 15
    build_thank_you_slide(prs)          # 16

    prs.save(str(OUTPUT_PATH))
    print(f"Presentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
